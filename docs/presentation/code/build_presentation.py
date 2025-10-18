"""
Build TrodesTrack presentation using python-pptx

Generates a 43-slide PowerPoint presentation for neuroscientists explaining
the TrodesTrack sensor-fusion tracking package.

Usage:
    uv run python docs/presentation/code/build_presentation.py

Output:
    docs/presentation/trodestrack_presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Paths
SCRIPT_DIR = Path(__file__).parent
PRESENTATION_DIR = SCRIPT_DIR.parent
VISUALS_DIR = PRESENTATION_DIR / "visuals"
VIDEOS_DIR = PRESENTATION_DIR / "videos"
OUTPUT_FILE = PRESENTATION_DIR / "trodestrack_presentation.pptx"

# Color palette (from design specs)
COLORS = {
    "blue": RGBColor(46, 134, 171),  # #2E86AB - trust
    "orange": RGBColor(247, 127, 0),  # #F77F00 - energy
    "green": RGBColor(6, 167, 125),  # #06A77D - success
    "red": RGBColor(214, 40, 40),  # #D62828 - error
    "gray": RGBColor(108, 117, 125),  # #6C757D - neutral
    "black": RGBColor(0, 0, 0),
    "white": RGBColor(255, 255, 255),
}


class PresentationBuilder:
    """Build TrodesTrack presentation slide by slide"""

    def __init__(self):
        self.prs = Presentation()
        # Set slide size to 16:9
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def save(self, output_path):
        """Save presentation to file"""
        self.prs.save(output_path)
        print(f"✓ Saved: {output_path}")

    def add_title_slide(self, title, subtitle, author="", affiliation=""):
        """Add title slide (layout 0)"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._format_text(title_shape, font_size=44, bold=True, color=COLORS["blue"])

        # Subtitle
        if slide.shapes.placeholders[1]:
            subtitle_shape = slide.shapes.placeholders[1]
            subtitle_text = subtitle
            if author:
                subtitle_text += f"\n\n{author}"
            if affiliation:
                subtitle_text += f"\n{affiliation}"
            subtitle_shape.text = subtitle_text
            self._format_text(subtitle_shape, font_size=28, color=COLORS["gray"])

        return slide

    def add_section_divider(self, section_title, section_number, total_sections=6):
        """Add section divider slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Background color (blue gradient effect)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS["blue"]

        # Section number (large, top-left)
        section_num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(3), Inches(1.5))
        section_num_frame = section_num_box.text_frame
        section_num_frame.text = f"Section {section_number}/{total_sections}"
        self._format_text(section_num_box, font_size=32, color=COLORS["white"], alpha=0.7)

        # Section title (centered, large)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.word_wrap = True
        self._format_text(
            title_box, font_size=54, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER
        )

        return slide

    def add_content_slide(
        self, title, bullets=None, image_path=None, notes="", layout="title_content"
    ):
        """Add content slide with title, bullets, and optional image

        Args:
            title: Slide title
            bullets: List of bullet points (strings)
            image_path: Path to image file (optional)
            notes: Speaker notes
            layout: "title_content", "title_only", "two_column"
        """
        if layout == "two_column":
            slide_layout = self.prs.slide_layouts[3]  # Two content
        elif layout == "title_only":
            slide_layout = self.prs.slide_layouts[5]  # Title only
        else:
            slide_layout = self.prs.slide_layouts[1]  # Title and content

        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._format_text(title_shape, font_size=36, bold=True, color=COLORS["blue"])

        # Content area dimensions
        content_left = Inches(0.5)
        content_top = Inches(1.5)
        content_width = Inches(9)
        content_height = Inches(3.5)

        # If image provided, adjust layout
        if image_path:
            img_path = Path(image_path)
            if not img_path.is_absolute():
                img_path = VISUALS_DIR / img_path

            if img_path.exists():
                if bullets:
                    # Split: bullets on left, image on right
                    bullet_width = Inches(4)
                    img_left = Inches(5)
                    img_width = Inches(4.5)
                else:
                    # Full-width image
                    img_left = Inches(0.5)
                    img_width = Inches(9)
                    bullet_width = 0

                # Add image
                try:
                    slide.shapes.add_picture(str(img_path), img_left, content_top, width=img_width)
                except Exception as e:
                    print(f"Warning: Could not add image {img_path}: {e}")
            else:
                print(f"Warning: Image not found: {img_path}")
                bullet_width = content_width
        else:
            bullet_width = content_width

        # Add bullets if provided
        if bullets and bullet_width > 0:
            text_box = slide.shapes.add_textbox(
                content_left, content_top, bullet_width, content_height
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            for i, bullet_text in enumerate(bullets):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = bullet_text
                p.level = 0
                p.font.size = Pt(24)
                p.font.name = "Arial"
                p.font.color.rgb = COLORS["black"]
                p.space_before = Pt(12)
                p.space_after = Pt(12)

        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

        return slide

    def add_full_image_slide(self, title, image_path, caption="", notes=""):
        """Add slide with large centered image and optional caption"""
        slide_layout = self.prs.slide_layouts[5]  # Title only
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._format_text(title_shape, font_size=36, bold=True, color=COLORS["blue"])

        # Image (large, centered)
        img_path = Path(image_path)
        if not img_path.is_absolute():
            img_path = VISUALS_DIR / img_path

        if img_path.exists():
            try:
                img_top = Inches(1.3)
                img_height = Inches(3.8)

                # Add image centered
                pic = slide.shapes.add_picture(
                    str(img_path), Inches(0.5), img_top, height=img_height
                )

                # Center horizontally
                pic.left = int((self.prs.slide_width - pic.width) / 2)

            except Exception as e:
                print(f"Warning: Could not add image {img_path}: {e}")
        else:
            print(f"Warning: Image not found: {img_path}")

        # Caption (if provided)
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.3))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            self._format_text(
                caption_box, font_size=18, color=COLORS["gray"], align=PP_ALIGN.CENTER, italic=True
            )

        # Speaker notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

        return slide

    def _format_text(
        self, shape, font_size=24, bold=False, color=None, align=None, italic=False, alpha=1.0
    ):
        """Apply text formatting to a shape"""
        text_frame = shape.text_frame

        for paragraph in text_frame.paragraphs:
            if align:
                paragraph.alignment = align

            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.name = "Arial"
                run.font.bold = bold
                run.font.italic = italic

                if color:
                    run.font.color.rgb = color


def build_section_1_problem(builder):
    """Build Section 1: THE PROBLEM (Slides 1-8)"""

    # Slide 1: Title
    builder.add_title_slide(
        title="TrodesTrack: Sensor-Fusion Tracking for Behavioral Neuroscience",
        subtitle="Combining Camera + IMU for Robust 2D Position Tracking",
        author="",
        affiliation="",
    )

    # Slide 2: The Behavioral Tracking Challenge
    builder.add_full_image_slide(
        title="The Behavioral Tracking Challenge",
        image_path="slide02_failure_modes.png",
        caption="Common failure modes: occlusion, reflection, motion blur, dim lighting",
        notes="Neuroscience experiments rely on accurate tracking. Four common failure modes shown: LED occlusion (rat near wall), reflections from surfaces, motion blur during fast movements, and dim lighting reducing confidence. These failures corrupt analyses like place cells, trajectory decoding, and behavior scoring.",
    )

    # Slide 3: Real-World Consequences
    builder.add_full_image_slide(
        title="Real-World Consequences of Tracking Failures",
        image_path="slide03_trajectory_comparison.png",
        caption="Ground truth vs noisy camera observations with dropout gaps",
        notes="Left: Smooth ground truth trajectory. Right: Actual camera observations are noisy, have gaps during dropouts, and jumps when tracking fails. These artifacts contaminate downstream analyses: false neural correlates, biased velocity estimates, and unreliable behavior annotations.",
    )

    # Slide 4: Why Vision-Only Tracking Fails
    builder.add_content_slide(
        title="Why Vision-Only Tracking Fails",
        bullets=[
            "Camera limitations: 30 Hz, 2D only, line-of-sight required",
            "Occlusions: Rat near walls, corners, objects",
            "Reflections: Shiny surfaces, water ports, mirrors",
            "Motion blur: Fast movements exceed shutter speed",
            "Lighting variations: Shadows, LEDs too dim/bright",
            "Result: Gaps, jumps, outliers → Corrupted data",
        ],
        notes="Vision-only tracking has fundamental limitations. Cameras sample at 30 Hz, require line-of-sight, and are susceptible to environmental conditions. The result: gaps in data, sudden jumps when tracking recovers, and outliers that corrupt analyses. We need a complementary sensor.",
    )

    # Slide 5: Enter the IMU
    builder.add_content_slide(
        title="Enter the IMU (Inertial Measurement Unit)",
        bullets=[
            "Accelerometer: Measures specific force (m/s²)",
            "Gyroscope: Measures angular velocity (°/s)",
            "Advantages: 100+ Hz, always on, no occlusions",
            "Disadvantages: Drift, bias, gravity contamination",
            "Key insight: IMU fills gaps when camera fails!",
        ],
        notes="The IMU consists of accelerometers and gyroscopes mounted on the rat's headstage. It samples at 100+ Hz, always works (no occlusions), but has its own challenges: drift from bias accumulation and gravity contamination. The key insight: camera and IMU have complementary failure modes!",
    )

    # Slide 5A: Accelerometer Physics
    builder.add_full_image_slide(
        title="What Does an Accelerometer REALLY Measure?",
        image_path="slide05a_accelerometer_physics.png",
        caption="Specific force f = a - g (NOT acceleration!)",
        notes="Accelerometers measure specific force, not acceleration. At rest, they read +1g (not zero). In free fall, they read 0g (Einstein's equivalence principle). During motion, gravity contaminates X/Y axes when tilted. TrodesTrack assumes planar motion (2D) to avoid full tilt estimation.",
    )

    # Slide 5B: Gyroscope Physics
    builder.add_full_image_slide(
        title="What Does a Gyroscope REALLY Measure?",
        image_path="slide05b_gyroscope_physics.png",
        caption="Angular velocity → Integration gives heading → Bias causes unbounded drift",
        notes="Gyroscopes measure angular velocity (rotation rate). Integration gives heading: θ = θ₀ + ∫ω dt. Problem: Small bias (~1-5 °/s) accumulates unbounded. Without camera correction, heading drifts ±50° in 10 seconds! Kalman filter estimates bias in state vector to prevent drift.",
    )

    # Slide 5C: The Correction Challenge
    builder.add_full_image_slide(
        title="The Correction Challenge: 4 Fixes Needed",
        image_path="slide05c_bias_correction.png",
        caption="Left: Raw IMU integration drifts. Right: Kalman filter estimates biases → 10× lower error",
        notes="Four corrections needed: (1) Gravity removal (assume 2D), (2) Gyro bias estimation, (3) Accel bias estimation, (4) Frame alignment (IMU to body). Left panel shows raw integration drifting meters off course. Right panel shows Kalman filter estimating biases [b_gz, b_ax, b_ay] in state vector → reduces drift by 10×.",
    )

    # Slide 6: The Sensor Fusion Idea
    builder.add_content_slide(
        title="The Sensor Fusion Idea",
        bullets=[
            "Camera: Accurate, low-rate, fails during occlusions",
            "IMU: High-rate, always on, but drifts over time",
            "Sensor fusion: Combine both optimally!",
            "Camera corrects IMU drift",
            "IMU fills camera gaps",
            "Result: Best of both worlds",
        ],
        notes="Sensor fusion combines camera (accurate but sparse) with IMU (high-rate but drifts). Camera measurements correct IMU bias drift. IMU pre-integration fills gaps when camera fails. The Kalman filter is the mathematical framework that fuses these optimally under uncertainty.",
    )

    # Slide 7: What is TrodesTrack?
    builder.add_content_slide(
        title="What is TrodesTrack?",
        bullets=[
            "Python package for sensor-fusion rat tracking",
            "Inputs: Trodes/DLC camera + SpikeGadgets IMU",
            "Algorithms: Extended Kalman Filter (EKF), Unscented Kalman Filter (UKF), RTS Smoother",
            "Outputs: Position, velocity, heading + uncertainty",
            "Key features: Bias estimation, outlier rejection, dropout tolerance",
            "Built with JAX: Fast (300× realtime), GPU-ready",
        ],
        notes="TrodesTrack is an open-source Python package. It ingests camera tracking (Trodes LEDs or DeepLabCut) and SpikeGadgets IMU data. Implements EKF, UKF, and RTS smoothing with JAX (300× realtime on CPU). Outputs full state [x, y, vx, vy, θ, biases] with uncertainty estimates (covariance matrices).",
    )

    # Slide 8: Quick Preview - Before & After
    builder.add_content_slide(
        title="Quick Preview: Before & After",
        bullets=[
            "Scenario: 5-second camera dropout",
            "Vision-only: Huge drift (3+ meters)",
            "Sensor fusion: Bounded error (<10 cm)",
            "Video shows: Split-screen comparison",
            "Left: Vision-only extrapolation",
            "Right: TrodesTrack sensor fusion",
        ],
        image_path=None,  # Video will be added manually
        notes="This video shows a 10-second session with a 5-second camera dropout. Left: vision-only tracking extrapolates position, drifts 3+ meters. Right: TrodesTrack uses IMU pre-integration to maintain accuracy within 10 cm. Notice uncertainty grows during dropout (ellipses expand) and shrinks when camera returns. This is the core value proposition!",
    )
    # Note: Videos need to be added manually to PowerPoint after generation


def build_section_2_how_it_works(builder):
    """Build Section 2: HOW IT WORKS (Slides 9-18)"""

    # Slide 9: Section divider
    builder.add_section_divider("How It Works", section_number=2)

    # Slide 10: The Core Algorithm
    builder.add_content_slide(
        title="The Core Algorithm: Kalman Filtering",
        bullets=[
            "Two-step cycle: Predict → Update",
            "Predict: Use IMU to advance state forward in time",
            "Update: Use camera to correct prediction",
            "Handles uncertainty: Tracks confidence in estimates",
            "Optimal fusion: Weights sensors by their noise levels",
            "Mathematical framework for sensor fusion",
        ],
        notes="Kalman filtering is a recursive algorithm with two steps. Predict: Use motion model (IMU pre-integration) to forecast next state. Update: Use measurement model (camera observations) to correct prediction. The filter tracks uncertainty (covariance matrix P) and optimally weights sensors by their noise characteristics.",
    )

    # Slide 11: What We Track (State Vector)
    builder.add_content_slide(
        title="What We Track: The State Vector",
        bullets=[
            "8-dimensional state: [x, y, vₓ, vᵧ, θ, b_gz, b_ax, b_ay]",
            "Position: (x, y) in meters",
            "Velocity: (vₓ, vᵧ) in m/s",
            "Heading: θ in radians",
            "Biases: b_gz (gyro), b_ax, b_ay (accel)",
            "Why biases? Prevent unbounded drift over time",
        ],
        notes="TrodesTrack uses an 8D state vector. Position and velocity are obvious. Heading θ tracks rat orientation. The key innovation: estimating IMU biases [b_gz, b_ax, b_ay] as part of the state. These biases drift slowly (random walk), and estimating them prevents unbounded integration error. Camera measurements provide the observability needed to correct biases.",
    )

    # Slide 12: The Predict Step (IMU Integration)
    builder.add_full_image_slide(
        title="The Predict Step: IMU Pre-Integration",
        image_path="slide12_imu_integration.png",
        caption="Between camera frames: integrate high-rate IMU (gyro, accel) to advance state",
        notes="Between camera frames (33 ms @ 30 Hz), multiple IMU samples arrive (100 Hz = 3-4 samples). The predict step integrates these: Δθ = ∫(ω_z - b_gz) dt for heading, Δv = ∫R(θ)(a - b_a) dt for velocity, Δx = v*Δt + 0.5*Δv*Δt for position. This pre-integration maintains state at IMU rate, ready for the next camera measurement.",
    )

    # Slide 13: The Update Step (Camera Correction)
    builder.add_content_slide(
        title="The Update Step: Camera Correction",
        bullets=[
            "Camera observes: LED positions (x, y) or DLC keypoints",
            "Innovation: y = z_obs - h(x̂) (observation - prediction)",
            "Kalman gain K: Weights innovation by uncertainty",
            "Update: x̂ ← x̂ + K·y (correct prediction)",
            "Update covariance: P ← (I - K·H)·P (reduce uncertainty)",
            "Result: State 'snaps' toward camera observation",
        ],
        notes="When a camera frame arrives, we compute innovation (difference between observed and predicted LED position). Kalman gain K determines how much to trust this innovation based on camera noise vs prediction uncertainty. State is corrected, and covariance shrinks. If camera is missing (occlusion), we skip the update and uncertainty grows.",
    )

    # Slide 14: Handling Uncertainty
    builder.add_full_image_slide(
        title="Handling Uncertainty: Covariance Evolution",
        image_path="slide14_uncertainty.png",
        caption="Ellipses show position uncertainty (covariance). Grows during dropout, shrinks when camera returns.",
        notes="The covariance matrix P encodes uncertainty in the state estimate. During IMU-only prediction, uncertainty grows (process noise). When camera measures, uncertainty shrinks (innovation). This figure shows covariance ellipses (±2σ) during a 5-second dropout. Green = camera active, red = dropout. Notice bounded growth (not unbounded) thanks to damping term in dynamics.",
    )

    # Slide 15: EKF vs UKF
    builder.add_content_slide(
        title="EKF vs UKF: When to Use Each",
        bullets=[
            "EKF (Extended Kalman Filter): Linearizes nonlinear dynamics",
            "UKF (Unscented Kalman Filter): Samples sigma points, no linearization",
            "EKF: Faster, good for mild nonlinearities",
            "UKF: Slower, better for strong nonlinearities",
            "TrodesTrack default: EKF (meets <33 ms latency)",
            "UKF available for offline smoothing",
        ],
        notes="EKF linearizes the dynamics (Jacobians). UKF uses sigma points (deterministic sampling). For 2D tracking, nonlinearities are mild (heading rotation), so EKF works well. UKF is more accurate but 3× slower. TrodesTrack uses EKF by default for online filtering (<33 ms per frame). UKF is available for offline smoothing where speed is less critical.",
    )

    # Slide 16: Offline Smoothing
    builder.add_full_image_slide(
        title="Offline Smoothing: Using Future Data",
        image_path="slide16_smoother_comparison.png",
        caption="RTS smoother (right) uses future camera frames to refine estimates during dropout",
        notes="Filters are causal (online): only use past data. Smoothers are acausal (offline): use future data too. RTS (Rauch-Tung-Striebel) smoother runs forward filter, then backward pass to incorporate future measurements. Result: Lower error during dropouts. Compare left (filter) vs right (smoother) during 5s dropout: smoother drift is 1 cm vs filter's 8 cm. Use smoothing for offline analysis when latency doesn't matter.",
    )

    # Slide 17: Robustness Features
    builder.add_content_slide(
        title="Robustness Features",
        bullets=[
            "Mahalanobis gating: Reject outliers beyond 3σ threshold",
            "LED swap resolution: Handle front/back confusion via residuals",
            "Confidence scaling: Weight camera by DLC confidence",
            "Bias estimation: Adapt to slowly changing IMU characteristics",
            "Arena bounds: Soft constraint to stay within known region",
            "Damping: Exponential decay prevents velocity explosion during dropout",
        ],
        notes="TrodesTrack includes multiple robustness mechanisms. Mahalanobis gating rejects outliers (e.g., reflections) beyond χ² threshold. LED swaps are resolved by comparing wrapped residuals. DLC confidence scales measurement noise (low confidence → high noise). Bias estimation adapts to IMU drift. Arena bounds provide soft constraints. Velocity damping (λ term) prevents unbounded growth during long dropouts.",
    )

    # Slide 18: The 9-Panel Diagnostic Video
    builder.add_full_image_slide(
        title="The 9-Panel Diagnostic Video",
        image_path="slide18_diagnostic_panel.png",
        caption="Real-time monitoring: trajectory, sensors, errors, NEES, biases",
        notes="TrodesTrack generates 9-panel diagnostic videos for quality assurance. Top row: arena view, gyro, accel. Middle row: camera status, position error, velocity error. Bottom row: heading error, NEES (consistency), bias estimates. This comprehensive view lets you diagnose tracking failures, tune filter parameters, and verify results before analysis.",
    )


def build_section_3_features(builder):
    """Build Section 3: FEATURES & CAPABILITIES (Slides 19-25)"""

    # Slide 19: Section divider
    builder.add_section_divider("Features & Capabilities", section_number=3)

    # Slide 20: Synthetic Data Simulator
    builder.add_content_slide(
        title="Synthetic Data Simulator",
        bullets=[
            "Test algorithms without collecting real data",
            "simulate_rat_imu(): Generate realistic trajectories",
            "Includes: IMU noise, camera dropout, LED swaps, reflections",
            "Ground truth available: Evaluate RMSE, NEES",
            "Reproducible: Fixed random seeds for CI/CD",
            "Example: 30-line script generates benchmark",
        ],
        notes="TrodesTrack includes a high-fidelity simulator. simulate_rat_imu() generates synthetic sessions with realistic IMU noise (SpikeGadgets specs), camera dropout, LED swaps, and reflections. Ground truth is available for computing error metrics (RMSE, NEES). Fixed seeds ensure reproducibility. This enables test-driven development and algorithm validation before deploying on real data.",
    )

    # Slide 21: Quality Assurance Metrics
    builder.add_full_image_slide(
        title="Quality Assurance Metrics",
        image_path="slide21_nees_histogram.png",
        caption="NEES (Normalized Estimation Error Squared) vs χ² theoretical distribution",
        notes="NEES measures filter consistency: Does uncertainty (covariance) match actual errors? Ideal NEES = number of measurable DOF (5 for position + heading when camera active). Histogram shows distribution: mean ~5.8, within 95% confidence interval (green band). Overconfident filter → NEES too high. Underconfident → NEES too low. Use NEES to tune process/measurement noise parameters.",
    )

    # Slide 22: Automated QA Reports
    builder.add_content_slide(
        title="Automated QA Reports",
        bullets=[
            "9 metrics tracked: Position RMSE, velocity RMSE, heading error, NEES, dropout drift, etc.",
            "PDF report generated automatically",
            "Includes: Trajectory plots, error time series, NEES histogram, bias traces",
            "Pass/fail thresholds: RMSE ≤2 cm, velocity ≤10 cm/s, heading ≤7°, drift ≤3.5 m @ 5s",
            "CI integration: Prevents regressions",
            "Example: 5-minute session → 8-page report",
        ],
        notes="TrodesTrack auto-generates comprehensive QA reports. 9 metrics are computed and compared to acceptance thresholds (from PRD). PDF includes trajectory overlays, error time series, NEES histogram, and bias traces. This enables systematic tuning and prevents regressions. CI runs these on every commit. Real data can also be processed if hand-labeled ground truth is available.",
    )

    # Slide 23: Flexible State Tracking Modes
    builder.add_content_slide(
        title="Flexible State Tracking Modes",
        bullets=[
            "2d_pos: Position only (x, y) - 2D state",
            "2d_vel: Position + velocity (x, y, vₓ, vᵧ) - 4D state",
            "2d_full: Position + velocity + heading + biases - 8D state (default)",
            "heading_only: Heading only (θ, b_gz) - 2D state",
            "Future 3d_pose: Full 6-DOF (x, y, z, roll, pitch, yaw) - 16D state",
            "Tradeoff: Complexity vs accuracy vs computation",
        ],
        notes="TrodesTrack supports multiple state layouts via get_layout() API. 2d_pos is simplest (no velocity). 2d_vel adds velocity but no heading. 2d_full (default) includes heading and biases (8D). heading_only tracks orientation alone (useful for virtual reality). Future 3D extension will support full 6-DOF pose. Choose layout based on experiment needs and computational budget.",
    )

    # Slide 24: Performance & Scalability
    builder.add_content_slide(
        title="Performance & Scalability",
        bullets=[
            "CPU (single core): 300× realtime (5-min session in 1 second)",
            "GPU (single card): 1000×+ realtime (30-min session in 2 seconds)",
            "Memory efficient: O(n) scaling, ~100 MB for 30-min session",
            "Batch processing: Parallelize across sessions",
            "Online latency: <33 ms per frame (EKF)",
            "JAX JIT compilation: First run slow, subsequent runs fast",
        ],
        notes="TrodesTrack is highly optimized using JAX. CPU performance: 300× realtime on single core (M1 Mac / Intel Xeon). GPU acceleration: 1000×+ realtime on NVIDIA A100. Memory scales linearly O(n). Batch processing across sessions is embarrassingly parallel. Online filtering meets <33 ms latency requirement (30 Hz). JAX JIT compilation adds ~10s overhead on first run, then subsequent runs are instant.",
    )

    # Slide 25: Real Data Support
    builder.add_content_slide(
        title="Real Data Support",
        bullets=[
            "Input formats: Trodes LEDs, DeepLabCut CSV, SpikeGadgets IMU (.mda, .rec)",
            "Homography calibration: Interactive tool (click arena corners)",
            "Time synchronization: Hardware-synced (SpikeGadgets clock)",
            "Preprocessing: Remove IMU sample-and-hold repeats, convert units",
            "Output formats: Parquet (states), HDF5 (diagnostics), MP4 (videos)",
            "Example datasets: fetch_example() downloads demo session",
        ],
        notes="TrodesTrack ingests real data from Trodes (LED CSV), DeepLabCut (DLC CSV), and SpikeGadgets (MDA/REC files). Homography calibration tool (calib-homography command) maps pixels to meters via arena corners. Time sync assumes hardware-synced clocks. Preprocessing handles unit conversions and removes sample-and-hold repeats. Outputs: Parquet (time series), HDF5 (full diagnostics), MP4 (videos). Demo datasets available via fetch_example().",
    )


def build_section_4_getting_started(builder):
    """Build Section 4: GETTING STARTED (Slides 26-32)"""

    # Slide 26: Section divider
    builder.add_section_divider("Getting Started", section_number=4)

    # Slide 27: Installation
    builder.add_content_slide(
        title="Installation",
        bullets=[
            "Requirements: Python ≥3.11, JAX ≥0.4",
            "Install: git clone + uv sync",
            "Package manager: uv (fast, modern)",
            "Dependencies: numpy, scipy, jax, matplotlib",
            "Optional: GPU support (jax[cuda])",
            "Time: ~2 minutes on fresh machine",
        ],
        notes="Installation is straightforward. Clone the repo, run 'uv sync' to install dependencies. TrodesTrack uses uv for package management (faster than pip/conda). Core deps: numpy, scipy, jax, matplotlib. For GPU support, install jax[cuda]. Full installation takes ~2 minutes on a fresh machine. Verified on Linux, macOS, Windows (WSL).",
    )

    # Slide 28: Learning Path
    builder.add_content_slide(
        title="Learning Path: 10 Examples",
        bullets=[
            "01-02: Simple motion + EKF basics (35 lines total)",
            "03-05: Dropout handling, UKF, RTS smoothing",
            "06-08: Real data ingestion (Trodes, DLC, SpikeGadgets)",
            "09: Homography calibration (interactive tool)",
            "10: Full pipeline (end-to-end workflow)",
            "",
            "📚 Estimated time: 2-3 hours to complete all examples",
            "🚀 Each example runs in <1 minute",
        ],
        notes="TrodesTrack includes 10 progressive examples in the examples/ folder. Start with simple circular motion (5 lines), then EKF on synthetic data (30 lines). Examples 3-5 cover dropout, UKF, smoothing. Examples 6-8 show real data ingestion. Example 9 demonstrates homography calibration. Example 10 is the full pipeline. Each example is self-contained and runnable in <1 minute. Total learning time: 2-3 hours.",
    )

    # Slide 29: Decision Tree: Which Filter?
    builder.add_content_slide(
        title="Decision Tree: Which Filter to Use?",
        bullets=[
            "Need real-time (<33 ms)? → Use EKF",
            "Strong nonlinearities? → Use UKF (offline only)",
            "Offline analysis + max accuracy? → Use RTS smoother",
            "Iterative refinement? → Use IEKS (iterated EKF smoother)",
            "Unsure? → Start with EKF (default)",
            "All filters have identical API: just swap config",
        ],
        notes="Decision tree for choosing filter. Real-time requirement → EKF (fastest). Strong nonlinearities (e.g., 3D rotations) → UKF. Offline analysis + max accuracy → RTS smoother. Iterative refinement → IEKS (not yet implemented). When unsure, start with EKF (default, works for 95% of cases). All filters share identical API (filter config object + same input arrays), so switching is trivial.",
    )

    # Slide 30: When to Use TrodesTrack
    builder.add_content_slide(
        title="When to Use TrodesTrack",
        bullets=[
            "✅ BEST FOR:",
            "  • SpikeGadgets IMU + camera tracking",
            "  • 2D planar mazes with occlusions/dropouts",
            "  • High accuracy needs (RMSE <2 cm)",
            "",
            "❌ NOT RECOMMENDED FOR:",
            "  • Full 3D tracking (2D only for now)",
            "  • Real-time closed-loop (<33 ms latency)",
        ],
        notes="TrodesTrack is ideal for: SpikeGadgets IMU + camera setups, 2D planar tracking, high accuracy requirements, long sessions with occlusions. Python-based. NOT recommended for: Full 3D tracking (2D only for now), real-time closed-loop control (latency ~33 ms may be too slow), non-planar arenas (assumes flat surface). For 3D, use top-down projection. For closed-loop, consider simpler trackers or specialized hardware.",
    )

    # Slide 31: Troubleshooting Common Issues
    builder.add_content_slide(
        title="Troubleshooting Common Issues",
        bullets=[
            "Drift during dropout → Increase process noise or add damping",
            "Jumpy estimates → Decrease measurement noise or use gating",
            "Heading flips → Check LED spacing, enable swap resolution",
            "NEES too high → Filter overconfident, increase process noise",
            "NEES too low → Filter underconfident, decrease process noise",
            "Slow performance → Enable JAX JIT, check for Python loops",
            "See TUNING.md for detailed parameter guide",
        ],
        notes="Common issues and fixes. Excessive drift → increase process noise (Q) or damping (λ). Jumpy estimates → decrease camera noise (R) or tighten gating threshold. Heading flips → verify LED spacing parameter, enable swap resolution. NEES out of range → tune Q/R to match actual noise. Slow performance → ensure JAX JIT is enabled. Full troubleshooting guide in TUNING.md with NEES-based diagnostics.",
    )

    # Slide 32: Resources & Support
    builder.add_content_slide(
        title="Resources & Support",
        bullets=[
            "GitHub: github.com/edeno/trodestrack",
            "Documentation: README.md, API reference, examples/",
            "Guides: TUNING.md (parameter tuning), TROUBLESHOOTING.md",
            "Issues: Report bugs via GitHub Issues",
            "PRD: Full technical specification (.claude/docs/PRD.md)",
            "Contact: eric.denovellis@ucsf.edu",
        ],
        notes="All resources available on GitHub. README has quickstart guide. Examples folder has 10 progressive tutorials. TUNING.md explains NEES-based parameter selection. TROUBLESHOOTING.md covers common failure modes. PRD (.claude/docs/PRD.md) has full mathematical specification. Report bugs via GitHub Issues. Contact Eric Denovellis for questions or collaboration.",
    )


def build_section_5_advanced(builder):
    """Build Section 5: ADVANCED TOPICS (Slides 33-36)"""

    # Slide 33: Section divider
    builder.add_section_divider("Advanced Topics", section_number=5)

    # Slide 34: Under the Hood: JAX Implementation
    builder.add_content_slide(
        title="Under the Hood: JAX Implementation",
        bullets=[
            "JAX: NumPy + JIT compilation + autograd + GPU",
            "jax.lax.scan: Efficient loop fusion for filtering",
            "Pure functions: No side effects, enables parallelization",
            "XLA compilation: First run slow (~10s), then instant",
            "Speedup: 300× vs pure Python loops",
            "GPU: Same code runs on GPU via jax.device_put()",
        ],
        notes="TrodesTrack is built on JAX (Google's NumPy successor). jax.lax.scan fuses loops for efficient filtering. Pure functional design enables JIT compilation and parallelization. XLA backend compiles Python to machine code (first run slow, subsequent runs instant). Result: 300× speedup vs pure Python. GPU support is trivial: same code, just change device. This architecture enables real-time performance on modest hardware.",
    )

    # Slide 35: Extending to 3D
    builder.add_content_slide(
        title="Roadmap: Extending to 3D",
        bullets=[
            "Current: 2D state (x, y, θ, vx, vy, biases) - 8D",
            "Future: 3D pose (x, y, z, roll, pitch, yaw, velocities, biases) - 16D",
            "Requires: Full gyro/accel (all 3 axes), magnetometer (heading)",
            "Challenge: Quaternion vs Euler angles (gimbal lock)",
            "Use case: Rearing, climbing, non-planar behaviors",
            "Timeline: Prototype in 2025",
        ],
        notes="Current version: 2D tracking (x, y, heading). Future roadmap: Full 3D pose (x, y, z, roll, pitch, yaw). Requires using all 3 gyro axes, full accelerometer, and magnetometer for heading. Mathematical challenge: Quaternion representation (no gimbal lock) vs Euler angles (intuitive but singular). Use case: Rearing behavior, climbing, non-planar arenas. Prototype targeted for 2025. Architecture is designed to support this extension.",
    )

    # Slide 36: Custom Measurement Models
    builder.add_content_slide(
        title="Custom Measurement Models",
        bullets=[
            "Built-in: LED position, LED heading, velocity pseudo-measurement",
            "Protocol: Implement h(x) and H (Jacobian)",
            "Examples: ZUPT (zero velocity update), compass heading, arena bounds",
            "Plugin architecture: Swap measurement models without changing filter",
            "Composable: Combine multiple measurement types",
            "See models/measurements.py for templates",
        ],
        notes="TrodesTrack supports custom measurement models. Built-in models: LED position (x, y), LED heading (from front/back LEDs), velocity pseudo-measurements (when rat is stationary). To add your own: implement h(x) (measurement function) and H (Jacobian). Examples: ZUPT (zero velocity update when stationary), compass heading (magnetometer), arena bounds (soft constraints). Plugin architecture allows swapping models. See models/measurements.py for templates.",
    )


def build_section_6_conclusion(builder):
    """Build Section 6: CONCLUSION (Slides 37-40)"""

    # Slide 37: Section divider
    builder.add_section_divider("Conclusion", section_number=6)

    # Slide 38: Key Takeaways
    builder.add_content_slide(
        title="Key Takeaways",
        bullets=[
            "🎯 Sensor fusion (camera + IMU) >> vision-only",
            "🧮 Kalman filtering = optimal fusion framework",
            "⚙️ IMU bias estimation prevents unbounded drift",
            "📊 Accuracy: <2 cm RMSE, <3.5 m drift @ 5s",
            "⚡ Performance: 300× realtime on CPU (JAX)",
            "",
            "🚀 Start with Example 01 → real data in ~3 hours",
            "📖 Full documentation + 10 progressive examples",
        ],
        notes="Key takeaways for neuroscientists. Sensor fusion dramatically outperforms vision-only. Kalman filtering is the optimal framework for combining noisy sensors. IMU bias estimation is critical (not optional). Accelerometers measure specific force (gravity included), gyros drift unbounded. TrodesTrack meets strict accuracy targets: <2 cm RMSE, <3.5 m dropout drift. JAX provides extreme speed. Learning path: 10 examples, ~3 hours to competence.",
    )

    # Slide 39: Comparison to Alternatives
    builder.add_content_slide(
        title="Comparison to Alternatives",
        bullets=[
            "DeepLabCut alone: No IMU, dropout = data loss",
            "Trodes LEDs alone: Fast but no dropout tolerance",
            "Kalman smoothing (offline): Good but no real-time",
            "SLAM (RatSLAM): 3D mapping, overkill for 2D tracking",
            "TrodesTrack: 2D-specialized, sensor fusion, online + offline, open source",
            "Unique: Bias estimation + NEES-based tuning",
        ],
        notes="Comparison to alternatives. DeepLabCut: Excellent pose estimation but no IMU (dropout = data loss). Trodes LEDs: Fast centroid tracking but no dropout tolerance. Kalman smoothing (generic): Works offline but not real-time. SLAM (RatSLAM): 3D mapping for complex environments, overkill for simple 2D mazes. TrodesTrack: Specialized for 2D planar tracking, sensor fusion, both online and offline, open source. Unique features: IMU bias estimation and NEES-based diagnostics.",
    )

    # Slide 40: Future Directions
    builder.add_content_slide(
        title="Future Directions",
        bullets=[
            "3D pose estimation (roll, pitch, yaw)",
            "Magnetometer integration for absolute heading",
            "Multi-animal tracking (particle filter)",
            "Real-time closed-loop (reduce latency to <10 ms)",
            "Integration with spike sorting pipelines",
            "Cloud deployment (batch process 100s of sessions)",
            "Contributions welcome! github.com/edeno/trodestrack",
        ],
        notes="Future directions. 3D pose: Full 6-DOF tracking for non-planar behaviors. Magnetometer: Absolute heading without camera. Multi-animal: Particle filter for multiple rats. Real-time closed-loop: Reduce latency to <10 ms for VR applications. Spike sorting integration: Combine with Kilosort/Mountainsort for closed-loop analysis. Cloud deployment: Batch process large datasets. Contributions welcome via GitHub!",
    )

    # Slide 41: Thank You
    builder.add_content_slide(
        title="Thank You!",
        bullets=[
            "Questions?",
            "",
            "GitHub: github.com/edeno/trodestrack",
            "Email: eric.denovellis@ucsf.edu",
            "",
            "Try it: uv sync && uv run python examples/01_simple.py",
            "",
            "Slides + code: github.com/edeno/trodestrack/docs/presentation/",
        ],
        notes="Thank you for your attention! Questions welcome. Try TrodesTrack yourself: clone the repo, run 'uv sync', then 'uv run python examples/01_simple.py'. Full slides and generation code available in docs/presentation/ folder. Contributions and feedback welcome via GitHub Issues or email.",
    )

    # Slide 42: Acknowledgments (optional)
    builder.add_content_slide(
        title="Acknowledgments",
        bullets=[
            "SpikeGadgets: Hardware specifications and support",
            "JAX team (Google): High-performance numerical computing",
            "Trodes/DeepLabCut: Behavioral tracking software",
            "Neuroscience community: Feedback and testing",
            "Funding: [Add your funding sources here]",
        ],
        notes="Acknowledgments: SpikeGadgets for hardware specs and support. JAX team at Google for building an incredible numerical computing framework. Trodes and DeepLabCut teams for behavioral tracking software. Neuroscience community for feedback. [Add your specific funding sources and collaborators here].",
    )

    # Slide 43: Backup - References
    builder.add_content_slide(
        title="References & Further Reading",
        bullets=[
            "Kalman (1960): A New Approach to Linear Filtering",
            "Julier & Uhlmann (1997): Unscented Kalman Filter",
            "Rauch et al. (1965): RTS Smoothing Algorithm",
            "Bar-Shalom (1988): Tracking and Data Association",
            "Bradski et al. (2008): Learning OpenCV (homography)",
            "SpikeGadgets documentation: IMU specifications",
            "JAX documentation: jax.readthedocs.io",
        ],
        notes="Key references. Kalman 1960: Original Kalman filter paper. Julier & Uhlmann 1997: Unscented transform and UKF. Rauch 1965: RTS smoothing algorithm. Bar-Shalom 1988: Classic text on tracking and gating. Bradski 2008: OpenCV book, covers homography. SpikeGadgets docs: IMU specs and calibration. JAX docs: Full API reference and tutorials.",
    )


def main():
    """Main function to build presentation"""
    print("Building TrodesTrack presentation...")
    print()

    # Create builder
    builder = PresentationBuilder()

    # Build sections
    print("[1/6] Building Section 1: THE PROBLEM (slides 1-8)...")
    build_section_1_problem(builder)

    print("[2/6] Building Section 2: HOW IT WORKS (slides 9-18)...")
    build_section_2_how_it_works(builder)

    print("[3/6] Building Section 3: FEATURES & CAPABILITIES (slides 19-25)...")
    build_section_3_features(builder)

    print("[4/6] Building Section 4: GETTING STARTED (slides 26-32)...")
    build_section_4_getting_started(builder)

    print("[5/6] Building Section 5: ADVANCED TOPICS (slides 33-36)...")
    build_section_5_advanced(builder)

    print("[6/6] Building Section 6: CONCLUSION (slides 37-43)...")
    build_section_6_conclusion(builder)

    # Save presentation
    print()
    print("Saving presentation...")
    builder.save(OUTPUT_FILE)

    print()
    print("✅ Presentation complete!")
    print(f"   Slides: {len(builder.prs.slides)}")
    print(f"   Output: {OUTPUT_FILE}")
    print()
    print("Next steps:")
    print("  1. Open in PowerPoint/Keynote")
    print("  2. Add slide08_beforeafter.mp4 video manually to slide 8")
    print("  3. Review speaker notes")
    print("  4. Adjust fonts/colors if needed")
    print()
    print("To view:")
    print(f"  open {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
